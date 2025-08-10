''' This is version 4.0 '''
# Activate the virtual environment "venv" using "venv\Scripts\activate" and install the dependencies.
# Dependencies in the file "dependencies.txt", run 
# pip install -r dependencies.txt 
# in the terminal.

# Section 1: Imports and Setup

import os, re, json, argparse, tempfile
from tqdm import tqdm
from pptx import Presentation
from PIL import Image
import pytesseract
import dateparser
import google.generativeai as genai
from pathlib import Path
from collections import defaultdict
import difflib

# Scalability Improvements for PowerPoint Inconsistency Detector
# Enhanced version maintaining original commenting style and adding detailed comments

## 1. OCR PERFORMANCE IMPROVEMENTS

# Add these imports at the top of your script
import concurrent.futures
import time
from functools import lru_cache
import hashlib

from dotenv import load_dotenv
load_dotenv()  # This loads the .env file

GEMINI_MODEL = "gemini-2.5-flash"
genai.configure(api_key=os.environ.get("my_api_key"))
model = genai.GenerativeModel(GEMINI_MODEL)

# Section 2: Enhanced Helper functions with caching and parallel processing

# Replace your existing OCR helper functions with these enhanced versions:

@lru_cache(maxsize=100)
def ocr_image_bytes_cached(image_hash, image_bytes):
    """This function takes image bytes and returns the OCR text with caching to avoid reprocessing same images"""
    with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
        # This above line creates a temporary file with a .png suffix
        # delete=False means the file will not be deleted when closed, allowing us to read it
        tmp.write(image_bytes) # Does the actual writing of the image bytes to the temporary file
        tmp_path = tmp.name # tmp_path is set to name of temporary file
    try:
        # Custom OCR config to speed up processing by limiting character set to business-relevant ones
        custom_config = r'--oem 3 --psm 6 -c tessedit_char_whitelist=0123456789ABCDEFGHIJKLMNOPQRSTUVWXYZabcdefghijklmnopqrstuvwxyz.,%-$€₹ '
        text = pytesseract.image_to_string(Image.open(tmp_path), config=custom_config)
        # "try" is for trying to open the temporary file and run OCR on it
    finally:
        os.remove(tmp_path) # deleting the temporary file after OCR is done
    return text

def ocr_image_bytes(image_bytes):
    """Enhanced OCR with caching - maintains same interface but prevents reprocessing identical images"""
    # Create hash for caching identical images - prevents redundant OCR operations
    image_hash = hashlib.md5(image_bytes).hexdigest()
    return ocr_image_bytes_cached(image_hash, image_bytes)

# 1. PPTX helper function (for text-based) - Enhanced with same logic as v3
def extract_text_from_pptx(pptx_path):
    """This is a helper function that extracts text and any embedded images, from a .pptx file"""
    prs = Presentation(pptx_path)
    # prs is the Presentation object that represents the .pptx file
    slides = []
    for i, slide in enumerate(prs.slides, start=1):
        text_chunks=[]
        images=[]
        for shape in slide.shapes:
            if hasattr(shape,"text") and shape.text.strip():
                # "hasattr" means that the shape has a "text" attribute, and "shape.text.strip()" means that the text is not empty
                # "strip()" is used to remove any leading or trailing whitespace from the text
                text_chunks.append(shape.text.strip())
            if hasattr(shape, "image"):
                try:
                    images.append(shape.image.blob)
                    # "shape.image.blob" is used to get the image bytes from the shape
                except Exception:
                    pass 
        slides.append({"slide":i , "text": " ".join(text_chunks) , "images":images})
        # "slide" is slide number, "text" is text extracted from the slide, "images" is a list of image bytes extracted from slide
    return slides

def extract_text_from_images_parallel(img_folder, max_workers=4):
    """This runs OCR on all images in the folder with parallel processing for faster performance"""
    slides = []
    img_files = sorted([f for f in os.listdir(img_folder) if f.lower().endswith((".png", ".jpg", ".jpeg"))])
    # Inclusion of .png, .jpg, when the assignment only has .jpeg files, makes the tool more generalized and scalable
    
    def process_image(args):
        """Helper function to process single image in parallel thread"""
        index, img_file = args
        img_path = os.path.join(img_folder, img_file)
        
        # Skip very small images as they rarely contain meaningful text but slow down processing
        img = Image.open(img_path)
        if img.size[0] * img.size[1] < 10000:
            return {"slide": index, "text": "", "images": []}
        
        text = pytesseract.image_to_string(img, config=r'--oem 3 --psm 6') # This is actual the OCR with optimization
        return {"slide": index, "text": text, "images": []}
    
    # Use ThreadPoolExecutor for parallel OCR processing - much faster than sequential
    with concurrent.futures.ThreadPoolExecutor(max_workers=max_workers) as executor:
        futures = [executor.submit(process_image, (i+1, img_file)) 
                  for i, img_file in enumerate(img_files)]
        
        # Collect results with progress tracking just like the original
        for future in tqdm(concurrent.futures.as_completed(futures), 
                          total=len(futures), desc="OCR images"):
            slides.append(future.result())
    
    # Sort by slide number since parallel processing can complete out of order
    return sorted(slides, key=lambda x: x["slide"])

# 4. Normalization function for slides - Enhanced with same detailed comments as v3
def normalize_slides(slide):
    """This function normalizes the slide text by removing extra spaces and newlines, and extracts numbers/dates"""
    text = slide["text"] or ""
    if not text.strip() and slide["images"]:
        for img in slide["images"]:
            text += " " + ocr_image_bytes(img)
    num_regex = re.compile(r"(?:₹|\$|€)?\s*[-+]?\d{1,3}(?:[,\d{3}]+)?(?:\.\d+)?%?")
    # This above line is a regex that matches numbers with optional currency symbols and percentages
    # Also, it finds separators like commas and periods
    # Like the currency symbols help it to identify the currency, and the percentage symbol helps it to identify the percentage numbers
    numbers = list(set(num_regex.findall(text)))
    # This above line finds all the numbers in the text and removes duplicates
    # The set() function is used to remove duplicates, and then it is converted back to a list to maintain the order of numbers.
    dates=[]
    for word in text.split():
        dp = dateparser.parse(word, settings={'PREFER_DAY_OF_MONTH':'first'})
        # This line basically parses the word as a date, and if it is a valid date, it returns the date object
        # 'first' there means that if the date is ambiguous, it will prefer the first day of the month as default
        if dp: #If the date parser returns a valid date, it will be added to the list
            dates.append(dp.isoformat())
    return {"slide":slide["slide"] , "text": text.strip() , "numbers":numbers , "dates":list(set(dates))}

## 2. API RATE LIMITING AND BATCHING

class GeminiRateLimiter:
    """Simple rate limiter to prevent hitting Gemini API limits during large processing"""
    def __init__(self, calls_per_minute=15):
        self.calls_per_minute = calls_per_minute
        self.call_times = [] # Track when each API call was made to enforce limits
    
    def wait_if_needed(self):
        """Check if we need to wait before making next API call to respect rate limits"""
        now = time.time()
        # Remove calls older than 1 minute as they don't count toward current rate limit
        self.call_times = [t for t in self.call_times if now - t < 60]
        
        # If we're at the limit, wait until we can make another call
        if len(self.call_times) >= self.calls_per_minute:
            sleep_time = 60 - (now - self.call_times[0]) + 1
            if sleep_time > 0:
                print(f"⏳ Rate limiting: waiting {sleep_time:.1f}s...")
                time.sleep(sleep_time)
        
        self.call_times.append(now) # Record this API call timestamp

# Initialize rate limiter globally like the model initialization in original code
gemini_limiter = GeminiRateLimiter(calls_per_minute=12)

# Section 4: Enhanced Gemini Analysis with rate limiting

def extract_json_from_text(t):
    """Extract the first JSON array substring from text (crude but practical)."""
    if not isinstance(t, str):
        return None
    start = t.find('[')
    end = t.rfind(']')
    if start != -1 and end != -1 and end > start:
        return t[start:end+1]
    return None

def call_gemini_with_retry(slide_objs, max_retries=3):
    """Enhanced Gemini call with rate limiting and retry logic to handle API failures gracefully"""
    # Wait if needed to respect API rate limits
    gemini_limiter.wait_if_needed()
    
    # Build slide texts for prompt - same format as original but with dynamic sizing
    slide_texts = []
    for s in slide_objs:
        # Reduce excerpt size for large batches to avoid token limits
        excerpt_size = 800 if len(slide_objs) > 5 else 1200
        ppt_excerpt = (s.get("pptx_text") or "")[:excerpt_size]
        ocr_excerpt = (s.get("ocr_text") or "")[:excerpt_size]
        slide_texts.append(f"SLIDE {s['slide']}:\nPPTX: {ppt_excerpt}\nOCR: {ocr_excerpt}")

    # Enhanced prompt focusing on critical issues to reduce processing time and improve accuracy
    prompt = f"""
    Analyze these {len(slide_objs)} presentation slides for critical inconsistencies only:
    
    Focus on HIGH-IMPACT issues:
    1. Major monetary conflicts (>20% difference)
    2. Significant time savings conflicts
    3. Mathematical errors in totals
    4. Direct contradictions
    
    Return JSON array with max 5 most critical issues.
    
    Slides:
    {chr(10).join(slide_texts)}
    """

    # Retry logic with exponential backoff to handle temporary API failures
    for attempt in range(max_retries):
        try:
            resp = model.generate_content(prompt)
            text = getattr(resp, "text", None) or str(resp)
            
            # Same JSON extraction logic as original call_gemini function
            json_blob = extract_json_from_text(text)
            if json_blob:
                return json.loads(json_blob)
            return json.loads(text) if text.strip().startswith('[') else []
            
        except Exception as e:
            if attempt == max_retries - 1:
                return [{"type": "llm_error", "error": str(e)}] # Same error format as original
            print(f"⚠️  Gemini attempt {attempt + 1} failed, retrying...")
            time.sleep(2 ** attempt) # Exponential backoff: 2s, 4s, 8s delays
    
    return []

## 3. MEMORY OPTIMIZATION

class SlideProcessor:
    """Memory-efficient slide processor that handles large presentations without loading everything at once"""
    
    def __init__(self, chunk_size=10):
        self.chunk_size = chunk_size # Number of slides to process in each batch
    
    def process_slides_in_chunks(self, pptx_path, images_path):
        """Process slides in chunks to manage memory usage - yields one chunk at a time instead of loading all"""
        # Get total slide counts without loading all content
        prs = Presentation(pptx_path)
        total_pptx_slides = len(prs.slides)
        
        img_files = [f for f in os.listdir(images_path) 
                    if f.lower().endswith((".png", ".jpg", ".jpeg"))]
        # Same file filtering as original - inclusion of .png, .jpg makes tool more generalized
        total_img_slides = len(img_files)
        
        max_slides = max(total_pptx_slides, total_img_slides)
        
        # Process in chunks instead of all at once - this is the key memory optimization
        for start_idx in range(0, max_slides, self.chunk_size):
            end_idx = min(start_idx + self.chunk_size, max_slides)
            yield self._process_chunk(pptx_path, images_path, start_idx, end_idx)
    
    def _process_chunk(self, pptx_path, images_path, start_idx, end_idx):
        """Process a chunk of slides using same logic as original but only for subset"""
        # Extract only needed slides from PPTX - not all slides like original
        prs = Presentation(pptx_path)
        # prs is the Presentation object that represents the .pptx file
        pptx_chunk = []
        
        for i in range(start_idx, min(end_idx, len(prs.slides))):
            slide = prs.slides[i]
            text_chunks = []
            images = []
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    # "hasattr" means that the shape has a "text" attribute, and "shape.text.strip()" means that the text is not empty
                    # "strip()" is used to remove any leading or trailing whitespace from the text
                    text_chunks.append(shape.text.strip())
                if hasattr(shape, "image"):
                    try:
                        images.append(shape.image.blob)
                        # "shape.image.blob" is used to get the image bytes from the shape
                    except Exception:
                        pass
            
            pptx_chunk.append({
                "slide": i + 1, 
                "text": " ".join(text_chunks), 
                "images": images
            })
            # "slide" is slide number, "text" is text extracted from the slide, "images" is a list of image bytes extracted from slide
        
        # Process corresponding image slides using same logic as original
        img_files = sorted([f for f in os.listdir(images_path) 
                           if f.lower().endswith((".png", ".jpg", ".jpeg"))])
        ocr_chunk = []
        
        for i in range(start_idx, min(end_idx, len(img_files))):
            img_file = img_files[i]
            img_path = os.path.join(images_path, img_file)
            text = pytesseract.image_to_string(Image.open(img_path)) # This is actual the OCR
            ocr_chunk.append({"slide": i + 1, "text": text, "images": []})
        
        # Normalize and combine slides using same functions as original
        pptx_norm = [normalize_slides(s) for s in pptx_chunk]
        ocr_norm = [normalize_slides(s) for s in ocr_chunk]
        
        # Merge per-slide using same logic as original main function
        combined_chunk = []
        chunk_size = max(len(pptx_norm), len(ocr_norm))
        
        for i in range(chunk_size):
            slide_num = start_idx + i + 1
            ppt = pptx_norm[i] if i < len(pptx_norm) else {
                "slide": slide_num, "text": "", "numbers": [], "dates": []
            }
            ocr = ocr_norm[i] if i < len(ocr_norm) else {
                "slide": slide_num, "text": "", "numbers": [], "dates": []
            }
            
            merged_text = (ppt["text"] + " " + ocr["text"]).strip()
            merged_numbers = list(set(ppt.get("numbers", []) + ocr.get("numbers", [])))
            merged_dates = list(set(ppt.get("dates", []) + ocr.get("dates", [])))
            combined_chunk.append({
                "slide": slide_num,
                "text": merged_text,
                "numbers": merged_numbers,
                "dates": merged_dates,
                "pptx_text": ppt["text"],
                "ocr_text": ocr["text"]
            })
        
        return combined_chunk

# Section 3: Enhanced Rule-Based Detections (from v3 with detailed comments)

def extract_value_from_string(s):
    """Extract numeric value from string, handling currency and units"""
    # Remove common prefixes/suffixes and extract number
    cleaned = re.sub(r'[₹$€,\s]', '', s)
    match = re.search(r'(\d+(?:\.\d+)?)', cleaned)
    return float(match.group(1)) if match else None

def normalize_time_unit(value, text_context):
    """Convert time values to consistent units (hours per month)"""
    text_lower = text_context.lower()
    
    if 'per year' in text_lower or 'annually' in text_lower:
        return value / 12  # converts yearly to monthly
    elif 'minutes' in text_lower or 'mins' in text_lower:
        return value / 60  # converts minutes to hours
    elif 'per month' in text_lower or 'monthly' in text_lower:
        return value  # already monthly hours
    else:
        return value  # Default

def detect_impact_value_conflicts(slides):
    """Detect conflicts in key impact metrics like savings amounts"""
    conflicts = []
    impact_patterns = [
        r'(\$\d+(?:\.\d+)?[MmKk]?)\s*(?:saved|impact|productivity)',
        r'(?:saved|impact|productivity).*?(\$\d+(?:\.\d+)?[MmKk]?)',
        r'(\$\d+(?:\.\d+)?[MmKk]?)\s*(?:in|of)?\s*(?:lost\s*)?productivity'
    ]
    
    found_values = defaultdict(list)
    
    for slide in slides:
        text = slide["text"]
        for pattern in impact_patterns:
            matches = re.findall(pattern, text, re.IGNORECASE)
            for match in matches:
                normalized_value = normalize_currency_value(match)
                if normalized_value:
                    found_values[normalized_value].append({
                        'slide': slide['slide'],
                        'raw_text': match,
                        'context': text[:200]
                    })
    
    # Check for conflicts and include slide information
    unique_values = list(found_values.keys())
    if len(unique_values) > 1:
        slides_involved = []
        for value_list in found_values.values():
            slides_involved.extend([item['slide'] for item in value_list])
        
        conflicts.append({
            "type": "impact_value_conflict",
            "severity": "high",
            "slide": min(slides_involved),  # Primary slide for compatibility
            "slides": sorted(list(set(slides_involved))),  # All slides involved
            "description": f"Found conflicting impact values: {unique_values}",
            "values": dict(found_values)
        })
    
    return conflicts

def normalize_currency_value(value_str):
    """Convert currency strings to comparable numeric values"""
    if not value_str:
        return None
    
    # Extract number and multiplier
    match = re.search(r'(\d+(?:\.\d+)?)([MmKk]?)', value_str.replace('$', ''))
    # character-wise meaning of the regex in above line: (?:\.\d+) means "match a decimal point followed by one or more digits"
    # "(\d+(?:\.\d+)?)" means "match one or more digits, followed optionally by a decimal point and more digits"
    # "([MmKk]?)" means "match an optional M or K (case insensitive)"
    # "(\d+(?:\.\d+)?)([MmKk]?)"" means "match a number with optional decimal followed by an optional M or K"
    if not match:
        return None
    
    number = float(match.group(1))
    multiplier = match.group(2).upper()
    
    if multiplier == 'M':
        return number * 1000000
    elif multiplier == 'K':
        return number * 1000
    else:
        return number

def detect_time_savings_conflicts(slides):
    """Detect conflicts in time savings metrics with unit normalization"""
    conflicts = []
    time_patterns = [
        r'(\d+(?:\.\d+)?)\s*(minutes?|mins?|hours?)\s*(?:saved\s*)?(?:per\s*slide)',
        r'(?:saved\s*)?(\d+(?:\.\d+)?)\s*(minutes?|mins?|hours?)\s*per\s*slide',
        r'(\d+(?:\.\d+)?)\s*(minutes?|mins?|hours?)\s*(?:per\s*consultant\s*monthly)',
        r'(\d+(?:\.\d+)?)\s*(?:hours?)\s*(?:saved\s*)?(?:per\s*consultant\s*monthly)'
    ]
    
    per_slide_savings = []
    monthly_savings = []
    
    for slide in slides:
        text = slide["text"]
        
        # Look for "per slide" savings
        slide_matches = re.findall(r'(\d+(?:\.\d+)?)\s*(minutes?|mins?|hours?)\s*(?:saved\s*)?(?:per\s*slide)', text, re.IGNORECASE)
        for value_str, unit in slide_matches:
            value = float(value_str)
            if 'min' in unit.lower():
                value = value  # Keep in minutes for per-slide comparison
            else:
                value = value * 60  # Convert hours to minutes
            
            per_slide_savings.append({
                'slide': slide['slide'],
                'value': value,
                'unit': 'minutes',
                'raw_text': f"{value_str} {unit}",
                'context': text
            })
        
        # Look for monthly savings
        monthly_matches = re.findall(r'(\d+(?:\.\d+)?)\s*hours?\s*(?:saved\s*)?(?:per\s*consultant\s*monthly)', text, re.IGNORECASE)
        for value_str in monthly_matches:
            monthly_savings.append({
                'slide': slide['slide'],
                'value': float(value_str),
                'unit': 'hours_per_month',
                'raw_text': f"{value_str} hours",
                'context': text
            })
    
    # Check for per-slide conflicts
    if len(set(item['value'] for item in per_slide_savings)) > 1:
        conflicts.append({
            "type": "time_per_slide_conflict",
            "severity": "medium",
            "description": "Found conflicting time savings per slide values",
            "values": per_slide_savings
        })
    
    return conflicts

def detect_sum_breakdown_conflicts(slides):
    """Detect when breakdown components don't sum to the claimed total"""
    conflicts = []
    
    for slide in slides:
        text = slide["text"]
        
        # Look for "X Hours Saved Per Consultant Monthly" pattern
        total_match = re.search(r'(\d+)\s*Hours?\s*Saved\s*Per\s*Consultant\s*Monthly', text, re.IGNORECASE)
        if total_match:
            claimed_total = int(total_match.group(1))
            
            # Look for individual time savings in the same slide
            individual_savings = re.findall(r'(\d+)\s*hours?\s*per\s*consultant\s*monthly', text, re.IGNORECASE)
            if len(individual_savings) > 1:  # More than just the total
                # Remove the total from the list
                individual_values = [int(x) for x in individual_savings if int(x) != claimed_total]
                
                if individual_values:
                    actual_sum = sum(individual_values)
                    if actual_sum != claimed_total:
                        conflicts.append({
                            "type": "sum_breakdown_mismatch",
                            "severity": "high",
                            "slide": slide['slide'],
                            "claimed_total": claimed_total,
                            "breakdown_sum": actual_sum,
                            "breakdown_values": individual_values,
                            "description": f"Claimed total ({claimed_total}) doesn't match sum of breakdown ({actual_sum})"
                        })
    
    return conflicts

def detect_unit_mixing_conflicts(slides):
    """Detect when same metrics are presented in different time units"""
    conflicts = []
    layout_optimization_values = []
    
    for slide in slides:
        text = slide["text"]
        
        if 'layout optimization' in text.lower():
            yearly_matches = re.findall(r'(\d+)\s*hours?\s*(?:saved\s*)?(?:per\s*consultant[/\s]*year)', text, re.IGNORECASE)
            for value in yearly_matches:
                layout_optimization_values.append({
                    'slide': slide['slide'],
                    'value': int(value),
                    'unit': 'per_year',
                    'normalized_monthly': int(value) / 12
                })
            
            monthly_matches = re.findall(r'(\d+)\s*hours?\s*(?:saved\s*)?(?:per\s*consultant\s*monthly)', text, re.IGNORECASE)
            for value in monthly_matches:
                layout_optimization_values.append({
                    'slide': slide['slide'],
                    'value': int(value),
                    'unit': 'per_month',
                    'normalized_monthly': int(value)
                })
    
    if len(layout_optimization_values) > 1:
        normalized_values = [item['normalized_monthly'] for item in layout_optimization_values]
        if len(set(normalized_values)) > 1:
            slides_involved = [item['slide'] for item in layout_optimization_values]
            
            conflicts.append({
                "type": "unit_mixing_confusion",
                "severity": "medium",
                "slide": min(slides_involved),  # Primary slide
                "slides": sorted(list(set(slides_involved))),  # All slides involved
                "metric": "Layout Optimization",
                "description": "Same metric presented in different time units causing potential confusion",
                "values": layout_optimization_values
            })
    
    return conflicts

# Enhanced numeric conflict detection with context awareness
def detect_contextual_numeric_conflicts(slides):
    """Enhanced version that considers context when detecting numeric conflicts"""
    metric_map = defaultdict(lambda: defaultdict(list))
    conflicts = []

    for slide in slides:
        text = slide["text"]
        
        patterns = [
            r'(?:saves an estimated|delivering|achieving)\s*(\d+(?:\.\d+)?)\s*(hours?|minutes?|mins?)',
            r'(\$\d+(?:\.\d+)?[MK]?)\s*(?:saved|impact)',
            r'(\d+(?:\.\d+)?x)\s*(?:faster|speed)',
            r'(\d+(?:\.\d+)?)\s*(hours?|minutes?)\s*(?:from|saved)',
        ]
        
        for pattern in patterns:
            matches = re.finditer(pattern, text, re.IGNORECASE)
            for match in matches:
                value = match.group(1)
                context = text[max(0, match.start()-50):match.end()+50].strip()
                
                category = categorize_metric(context)
                if category:
                    metric_key = f"{category}:{value}"
                    metric_map[metric_key][value].append({
                        'slide': slide['slide'],
                        'context': context,
                        'full_match': match.group(0)
                    })

    # Only report conflicts within same categories
    for metric_key, values in metric_map.items():
        if len(values) > 1:
            category = metric_key.split(':')[0]
            slides_involved = []
            for value_list in values.values():
                slides_involved.extend([item['slide'] for item in value_list])
            
            conflicts.append({
                "type": "contextual_numeric_conflict",
                "severity": "medium",
                "slide": min(slides_involved),
                "slides": sorted(list(set(slides_involved))),
                "category": category,
                "metric": metric_key,
                "description": f"Conflicting {category} values found across slides",
                "values": dict(values)
            })
    
    return conflicts

def categorize_metric(context):
    """Categorize a metric based on its context"""
    context_lower = context.lower()
    
    if any(word in context_lower for word in ['hours saved', 'minutes saved', 'time saved']):
        return 'time_savings'
    elif any(word in context_lower for word in ['$', 'productivity', 'impact']):
        return 'impact'
    elif any(word in context_lower for word in ['faster', 'speed', 'efficiency']):
        return 'efficiency'
    else:
        return 'other'

# 2. Detecting percentage sum issues (keeping original with detailed comments)
def detect_percent_sum_issues(slides, tolerance=2.0):
    """Checks if percentages in a slide sum to ~100%."""
    issues = []
    for s in slides:
        percents = re.findall(r"[-+]?\d{1,3}(?:\.\d+)?%", s["text"])
        if len(percents) >= 2:
            total = sum(float(p.strip("%")) for p in percents)
            if abs(total - 100.0) > tolerance:
                issues.append({"type": "percent_total_mismatch", "slide": s["slide"], "found_total": total, "details": percents})
    return issues

def get_default_paths():
    """Helper function to find default paths for PPTX and images folder"""
    base_dir = Path(__file__).resolve().parent  # script is in NOOGATASSIGNMENT/
    # the images are expected under NoogatAssignment/NoogatAssignment/
    images_dir = base_dir / "NoogatAssignment" / "NoogatAssignment"
    pptx_files = list(base_dir.rglob("*.pptx"))
    pptx_file = pptx_files[0] if pptx_files else None

    return str(pptx_file) if pptx_file else None, str(images_dir)

## 4. ENHANCED MAIN FUNCTION

def main_enhanced():
    """Enhanced main function with descriptive output"""
    default_pptx, default_images = get_default_paths()

    print("\n--- ENHANCED SCALABLE PowerPoint Inconsistency Detector ---")
    print(f"Default PPTX:   {default_pptx if default_pptx else 'Not found'}")
    print(f"Default Images: {default_images if os.path.isdir(default_images) else 'Not found'}")
    print("----------------------------------------------------------------")

    pptx_path = input(f"Enter PPTX path [{default_pptx}]: ").strip() or default_pptx
    images_path = input(f"Enter images folder [{default_images}]: ").strip() or default_images

    if not pptx_path or not os.path.isfile(pptx_path):
        print(f"❌ PPTX file not found: {pptx_path}")
        return
    if not images_path or not os.path.isdir(images_path):
        print(f"❌ Images folder not found: {images_path}")
        return

    print("\n✅ Paths confirmed:")
    print(f"PPTX:   {pptx_path}")
    print(f"Images: {images_path}")

    # Check presentation size
    prs = Presentation(pptx_path)
    slide_count = len(prs.slides)
    
    if slide_count > 50:
        print(f"⚠️  Large presentation detected ({slide_count} slides)")
        print("   Enabling chunked processing and parallel OCR...")
        use_chunked = True
        ocr_workers = min(8, os.cpu_count())
    else:
        use_chunked = False
        ocr_workers = 4

    print(f"\n✅ Processing {slide_count} slides...")
    print(f"OCR workers: {ocr_workers}")
    print(f"Chunked processing: {'Yes' if use_chunked else 'No'}")

    if use_chunked:
        # Use chunked processing for large presentations
        processor = SlideProcessor(chunk_size=15)
        all_rule_issues = []
        all_llm_issues = []
        total_slides_processed = 0
        
        for chunk in processor.process_slides_in_chunks(pptx_path, images_path):
            print(f"Processing chunk: slides {chunk[0]['slide']}-{chunk[-1]['slide']}")
            
            chunk_rule_issues = []
            chunk_rule_issues.extend(detect_percent_sum_issues(chunk))
            chunk_rule_issues.extend(detect_impact_value_conflicts(chunk))
            chunk_rule_issues.extend(detect_time_savings_conflicts(chunk))
            chunk_rule_issues.extend(detect_sum_breakdown_conflicts(chunk))
            chunk_rule_issues.extend(detect_unit_mixing_conflicts(chunk))
            chunk_rule_issues.extend(detect_contextual_numeric_conflicts(chunk))
            
            all_rule_issues.extend(chunk_rule_issues)
            
            if os.environ.get("my_api_key"):
                chunk_llm_issues = call_gemini_with_retry(chunk)
                if isinstance(chunk_llm_issues, list):
                    all_llm_issues.extend(chunk_llm_issues)
            
            total_slides_processed += len(chunk)
            print(f"✓ Processed {total_slides_processed}/{slide_count} slides")
        
        combined_slides = []
        rule_issues = all_rule_issues
        llm_issues = all_llm_issues
        
    else:
        # Standard processing for smaller presentations
        print("\n[1/6] Extracting text from PPTX slides...")
        raw_pptx_slides = extract_text_from_pptx(pptx_path)

        print("[2/6] Running parallel OCR on slide images...")
        raw_image_slides = extract_text_from_images_parallel(images_path, max_workers=ocr_workers)

        print("[3/6] Normalizing slides (extracting numbers/dates)...")
        pptx_norm = [normalize_slides(s) for s in tqdm(raw_pptx_slides, desc="Normalizing PPTX")]
        ocr_norm = [normalize_slides(s) for s in tqdm(raw_image_slides, desc="Normalizing OCR")]

        # Merge per-slide
        max_slides = max(len(pptx_norm), len(ocr_norm))
        combined_slides = []
        for i in range(max_slides):
            ppt = pptx_norm[i] if i < len(pptx_norm) else {"slide": i+1, "text": "", "numbers": [], "dates": []}
            ocr = ocr_norm[i] if i < len(ocr_norm) else {"slide": i+1, "text": "", "numbers": [], "dates": []}
            merged_text = (ppt["text"] + " " + ocr["text"]).strip()
            merged_numbers = list(set(ppt.get("numbers", []) + ocr.get("numbers", [])))
            merged_dates = list(set(ppt.get("dates", []) + ocr.get("dates", [])))
            combined_slides.append({
                "slide": i+1,
                "text": merged_text,
                "numbers": merged_numbers,
                "dates": merged_dates,
                "pptx_text": ppt["text"],
                "ocr_text": ocr["text"]
            })

        print("[4/6] Running enhanced rule-based detectors...")
        rule_issues = []
        
        rule_issues.extend(detect_percent_sum_issues(combined_slides))
        rule_issues.extend(detect_impact_value_conflicts(combined_slides))
        rule_issues.extend(detect_time_savings_conflicts(combined_slides))
        rule_issues.extend(detect_sum_breakdown_conflicts(combined_slides))
        rule_issues.extend(detect_unit_mixing_conflicts(combined_slides))
        rule_issues.extend(detect_contextual_numeric_conflicts(combined_slides))

        # LLM analysis
        llm_issues = []
        api_key_present = os.environ.get("my_api_key") is not None
        
        if api_key_present:
            print("[5/6] Running enhanced Gemini deep checks (batched)...")
            batch_size = 6
            for i in range(0, len(combined_slides), batch_size):
                batch = combined_slides[i:i+batch_size]
                res = call_gemini_with_retry(batch)
                if isinstance(res, list):
                    llm_issues.extend(res)
        else:
            print("⚠️  No Gemini API key found in environment variable 'my_api_key'. Skipping deep LLM checks.")

    # Enhanced summary generation
    print("[6/6] Generating enhanced report...")
    
    # Categorize issues by severity
    high_priority = [issue for issue in rule_issues if issue.get('severity') == 'high']
    medium_priority = [issue for issue in rule_issues if issue.get('severity') == 'medium']
    low_priority = [issue for issue in rule_issues if issue.get('severity') not in ['high', 'medium']]

    # Create lightweight output for large presentations
    slides_summary = [{"slide": i+1, "processed": True} for i in range(slide_count)]
    
    out = {
        "summary": {
            "slides_processed": slide_count,
            "total_issues": len(rule_issues) + len(llm_issues),
            "rule_issues": len(rule_issues),
            "llm_issues": len(llm_issues),
            "high_priority_issues": len(high_priority),
            "medium_priority_issues": len(medium_priority),
            "low_priority_issues": len(low_priority),
            "large_presentation": slide_count > 50
        },
        "rule_issues": rule_issues,
        "llm_issues": llm_issues,
        "slides_summary": slides_summary
    }
    
    # Include full slide data for smaller presentations only
    if not use_chunked and combined_slides:
        out["slides"] = [
            {"slide": s["slide"], "pptx_text": s["pptx_text"], "ocr_text": s["ocr_text"], "numbers": s["numbers"], "dates": s["dates"]}
            for s in combined_slides
        ]

    # Save output
    out_path = Path(__file__).resolve().parent / "inconsistencies_enhanced.json"
    with open(out_path, "w", encoding="utf-8") as f:
        json.dump(out, f, indent=2, ensure_ascii=False)

    # V3.0 STYLE DESCRIPTIVE TERMINAL OUTPUT
    print("\n=== Enhanced Analysis Complete ===")
    print(f"Slides processed: {out['summary']['slides_processed']}")
    print(f"Total issues found: {out['summary']['total_issues']}")
    print(f"  - High priority: {out['summary']['high_priority_issues']}")
    print(f"  - Medium priority: {out['summary']['medium_priority_issues']}")  
    print(f"  - Low priority: {out['summary']['low_priority_issues']}")
    print(f"Rule-based issues: {out['summary']['rule_issues']}")
    print(f"LLM issues: {out['summary']['llm_issues']}")
    print(f"Report saved to: {out_path}")

    # DETAILED ISSUE BREAKDOWN - V3.0 STYLE
    if high_priority:
        print("\n🚨 HIGH PRIORITY ISSUES DETECTED:")
        for issue in high_priority:
            issue_type = issue.get('type', 'Unknown').replace('_', ' ').title()
            description = issue.get('description', 'No description')
            
            # Handle slide number display properly
            if issue.get('slides'):
                slide_info = f"Slides {', '.join(map(str, issue['slides']))}"
            elif issue.get('slide'):
                slide_info = f"Slide {issue['slide']}"
            else:
                slide_info = "Multiple slides"
            
            print(f"  - {issue_type} ({slide_info}): {description}")
            
            # Add detailed breakdown for specific issue types
            if issue.get('type') == 'sum_breakdown_mismatch':
                claimed = issue.get('claimed_total', 'Unknown')
                actual = issue.get('breakdown_sum', 'Unknown')
                breakdown = issue.get('breakdown_values', [])
                print(f"    → Claimed Total: {claimed}")
                print(f"    → Actual Sum: {actual}")
                print(f"    → Breakdown: {breakdown}")
            
            if issue.get('values'):
                print(f"    → Evidence: {list(issue['values'].keys())}")
    
    if medium_priority:
        print("\n⚠️  MEDIUM PRIORITY ISSUES:")
        for issue in medium_priority:
            issue_type = issue.get('type', 'Unknown').replace('_', ' ').title()
            description = issue.get('description', 'No description')
            
            if issue.get('slides'):
                slide_info = f"Slides {', '.join(map(str, issue['slides']))}"
            elif issue.get('slide'):
                slide_info = f"Slide {issue['slide']}"
            else:
                slide_info = "Multiple slides"
                
            print(f"  - {issue_type} ({slide_info}): {description}")
            
            # Add specific details for time conflicts
            if issue.get('type') == 'time_per_slide_conflict' and issue.get('values'):
                print(f"    → Conflicting Values:")
                for val in issue['values']:
                    print(f"      • Slide {val['slide']}: {val['raw_text']}")
            
            # Add details for unit mixing
            if issue.get('type') == 'unit_mixing_confusion' and issue.get('values'):
                print(f"    → Mixed Units Found:")
                for val in issue['values']:
                    print(f"      • Slide {val['slide']}: {val['value']} {val['unit']}")
    
    # Enhanced LLM issues display - v3.0 style with more detail
    if llm_issues:
        valid_llm_issues = [issue for issue in llm_issues 
                           if issue.get('issue_type') != 'llm_error' and issue.get('type') != 'llm_error']
        
        if valid_llm_issues:
            print("\n🔍 AI-DETECTED ISSUES:")
            for i, issue in enumerate(valid_llm_issues, 1):
                # Handle different possible field names from LLM
                issue_type = (issue.get('issue_type') or 
                             issue.get('type') or 
                             'AI Analysis').replace('_', ' ').title()
                
                description = (issue.get('description') or 
                              issue.get('details') or 
                              str(issue)[:200])
                
                slides_involved = (issue.get('slides') or 
                                 issue.get('slide_numbers') or 
                                 [issue.get('slide')] if issue.get('slide') else [])
                
                # Clean up slides list
                if slides_involved:
                    slides_involved = [s for s in slides_involved if s is not None]
                
                if slides_involved:
                    slide_info = f"Slides {', '.join(map(str, slides_involved))}"
                    print(f"  {i}. {issue_type} ({slide_info}):")
                else:
                    print(f"  {i}. {issue_type}:")
                
                # Format description with proper wrapping
                import textwrap
                wrapped_desc = textwrap.fill(description, width=80, 
                                           initial_indent="     ", subsequent_indent="     ")
                print(wrapped_desc)
                
                # Add evidence if available
                if issue.get('evidence'):
                    print("     → Evidence:")
                    for evidence in issue['evidence'][:2]:  # Show first 2 pieces of evidence
                        wrapped_evidence = textwrap.fill(f"• {evidence}", width=75,
                                                       initial_indent="       ", subsequent_indent="         ")
                        print(wrapped_evidence)
                print()  # Add spacing between issues
    
    # Final summary
    if not rule_issues and not llm_issues:
        print("\n✅ NO INCONSISTENCIES DETECTED - All slides appear consistent!")
        print("   The presentation maintains logical consistency across all analyzed content.")
    else:
        print(f"\n📊 ANALYSIS SUMMARY:")
        print(f"   • Total Issues Found: {len(rule_issues) + len(llm_issues)}")
        print(f"   • Critical Issues: {len(high_priority)} (require immediate attention)")
        print(f"   • Warning Issues: {len(medium_priority)} (should be reviewed)")
        print(f"   • Rule-based Detection: {len(rule_issues)} issues")
        print(f"   • AI Analysis: {len(valid_llm_issues) if 'valid_llm_issues' in locals() else len(llm_issues)} issues")
    
    print(f"\n📄 Detailed report saved to: {out_path}")
    print("   This JSON file contains complete analysis data and evidence for all detected issues.")

if __name__ == "__main__":
    main_enhanced()