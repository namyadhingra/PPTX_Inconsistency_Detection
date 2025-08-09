'''This is the Third version'''
# Activate the virtual environment "venv" using "venv\Scripts\activate" and install the dependencies.
# Dependencies in the file "dependencies.txt", run 
# pip install -r dependencies.txt 
# in the terminal.

# Section 1: Imports and Setup

import os, re, json, argparse, tempfile
from tqdm import tqdm
from pptx import Presentation
from PIL import Image
import pytesseract
import dateparser
import google.generativeai as genai
from pathlib import Path
from collections import defaultdict
import difflib

from dotenv import load_dotenv
load_dotenv()  # This loads the .env file

GEMINI_MODEL = "gemini-2.5-flash"
genai.configure(api_key=os.environ.get("my_api_key"))
model = genai.GenerativeModel(GEMINI_MODEL)

# Section 2: Helper functions

# 1. PPTX helper function (for text-based)
def extract_text_from_pptx(pptx_path):
    """This is a helper function that extracts text and any embedded images, from a .pptx file"""
    prs = Presentation(pptx_path)
    # prs is the Presentation object that represents the .pptx file
    slides = []
    for i, slide in enumerate(prs.slides, start=1):
        text_chunks=[]
        images=[]
        for shape in slide.shapes:
            if hasattr(shape,"text") and shape.text.strip():
                # "hasattr" means that the shape has a "text" attribute, and "shape.text.strip()" means that the text is not empty
                # "strip()" is used to remove any leading or trailing whitespace from the text
                text_chunks.append(shape.text.strip())
            if hasattr(shape, "image"):
                try:
                    images.append(shape.image.blob)
                    # "shape.image.blob" is used to get the image bytes from the shape
                except Exception:
                    pass 
        slides.append({"slide":i , "text": " ".join(text_chunks) , "images":images})
        # "slide" is slide number, "text" is text extracted from the slide, "images" is a list of image bytes extracted from slide
    return slides

# 2. OCR helper function (for image-based)
def extract_text_from_images( img_folder):
    """This runs OCR on all images in the folder specified as "img_folder" """
    slides= []
    img_files = sorted( [f for f in os.listdir(img_folder) if f.lower().endswith((".png", ".jpg", ".jpeg"))])
    # Inclusion of .png, .jpg, when the assignment only has .jpeg files, makes the tool more generalized and scalable
    for index, img_file in enumerate(tqdm(img_files, desc="OCR images"), start=1):
        img_path = os.path.join(img_folder, img_file)
        text = pytesseract.image_to_string(Image.open(img_path)) #This is actual the OCR
        slides.append({"slide":index , "text":text , "images":[]})
    return slides

# 3. OCR for image bytes (image bytes = raw bytes)
def ocr_image_bytes (image_bytes):
    """This function takes image bytes and returns the OCR text"""
    with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
        #This above line creates a temporary file with a .png suffix
        # delete=False means the file will not be deleted when closed, allowing us to read it
        tmp.write(image_bytes) #Does the actual writing of the image bytes to the temporary file
        tmp_path = tmp.name #tmp_path is set to name of temporary file
    try:
        text = pytesseract.image_to_string(Image.open(tmp_path))
        # "try" is for trying to open the temporary file and run OCR on it
    finally:
        os.remove(tmp_path) #deleting the temporary file after OCR is done
    return text

# 4. Normalization function for slides
def normalize_slides(slide):
    """This function normalizes the slide text by removing extra spaces and newlines, and extracts numbers/dates"""
    text = slide["text"] or ""
    if not text.strip() and slide["images"]:
        for img in slide["images"]:
            text += " " + ocr_image_bytes(img)
    num_regex = re.compile(r"(?:₹|\$|€)?\s*[-+]?\d{1,3}(?:[,\d{3}]+)?(?:\.\d+)?%?")
    # This above line is a regex that matches numbers with optional currency symbols and percentages
    # Also, it finds separators like commas and periods
    # Like the currency symbols help it to identify the currency, and the percentage symbol helps it to identify the percentage numbers
    numbers = list(set(num_regex.findall(text)))
    # This above line finds all the numbers in the text and removes duplicates
    # The set() function is used to remove duplicates, and then it is converted back to a list to maintain the order of numbers.
    dates=[]
    for word in text.split():
        dp = dateparser.parse(word, settings={'PREFER_DAY_OF_MONTH':'first'})
        # This line basically parses the word as a date, and if it is a valid date, it returns the date object
        # 'first' there means that if the date is ambiguous, it will prefer the first day of the month as default
        if dp: #If the date parser returns a valid date, it will be added to the list
            dates.append(dp.isoformat())
    return {"slide":slide["slide"] , "text": text.strip() , "numbers":numbers , "dates":list(set(dates))}


# Section 3: Enhanced Rule-Based Detections

def extract_value_from_string(s):
    """Extract numeric value from string, handling currency and units"""
    # Remove common prefixes/suffixes and extract number
    cleaned = re.sub(r'[₹$€,\s]', '', s)
    match = re.search(r'(\d+(?:\.\d+)?)', cleaned)
    return float(match.group(1)) if match else None

def normalize_time_unit(value, text_context):
    """Convert time values to consistent units (hours per month)"""
    text_lower = text_context.lower()
    
    if 'per year' in text_lower or 'annually' in text_lower:
        return value / 12  # converts yearly to monthly
    elif 'minutes' in text_lower or 'mins' in text_lower:
        return value / 60  # converts minutes to hours
    elif 'per month' in text_lower or 'monthly' in text_lower:
        return value  # already monthly hours
    else:
        return value  # Default

def detect_impact_value_conflicts(slides):
    """Detect conflicts in key impact metrics like savings amounts"""
    conflicts = []
    impact_patterns = [
        r'(\$\d+(?:\.\d+)?[MmKk]?)\s*(?:saved|impact|productivity)',
        # Regex - matches a monetary value followed by keywords indicating impact
        r'(?:saved|impact|productivity).*?(\$\d+(?:\.\d+)?[MmKk]?)',
        # character to character description of above line:
        # "\$\d+(?:\.\d+" means a monetary value with optional decimal and million/billion suffix

        r'(\$\d+(?:\.\d+)?[MmKk]?)\s*(?:in|of)?\s*(?:lost\s*)?productivity'
        # Regex - matches a monetary value followed by optional keywords indicating lost productivity
    ]
    
    found_values = defaultdict(list)
    
    for slide in slides:
        text = slide["text"]
        for pattern in impact_patterns:
            matches = re.findall(pattern, text, re.IGNORECASE)
            for match in matches:
                # Normalize the value (convert M to millions, etc.)
                normalized_value = normalize_currency_value(match)
                if normalized_value:
                    found_values[normalized_value].append({
                        'slide': slide['slide'],
                        'raw_text': match,
                        'context': text[:200]
                    })
    
    # Check for conflicts
    unique_values = list(found_values.keys())
    if len(unique_values) > 1:
        conflicts.append({
            "type": "impact_value_conflict",
            "severity": "high",
            "description": f"Found conflicting impact values: {unique_values}",
            "values": dict(found_values)
        })
    
    return conflicts

def normalize_currency_value(value_str):
    """Convert currency strings to comparable numeric values"""
    if not value_str:
        return None
    
    # Extract number and multiplier
    match = re.search(r'(\d+(?:\.\d+)?)([MmKk]?)', value_str.replace('$', ''))
    if not match:
        return None
    
    number = float(match.group(1))
    multiplier = match.group(2).upper()
    
    if multiplier == 'M':
        return number * 1000000
    elif multiplier == 'K':
        return number * 1000
    else:
        return number

def detect_time_savings_conflicts(slides):
    """Detect conflicts in time savings metrics with unit normalization"""
    conflicts = []
    time_patterns = [
        r'(\d+(?:\.\d+)?)\s*(minutes?|mins?|hours?)\s*(?:saved\s*)?(?:per\s*slide)',
        r'(?:saved\s*)?(\d+(?:\.\d+)?)\s*(minutes?|mins?|hours?)\s*per\s*slide',
        r'(\d+(?:\.\d+)?)\s*(minutes?|mins?|hours?)\s*(?:per\s*consultant\s*monthly)',
        r'(\d+(?:\.\d+)?)\s*(?:hours?)\s*(?:saved\s*)?(?:per\s*consultant\s*monthly)'
    ]
    
    per_slide_savings = []
    monthly_savings = []
    
    for slide in slides:
        text = slide["text"]
        
        # Look for "per slide" savings
        slide_matches = re.findall(r'(\d+(?:\.\d+)?)\s*(minutes?|mins?|hours?)\s*(?:saved\s*)?(?:per\s*slide)', text, re.IGNORECASE)
        for value_str, unit in slide_matches:
            value = float(value_str)
            if 'min' in unit.lower():
                value = value  # Keep in minutes for per-slide comparison
            else:
                value = value * 60  # Convert hours to minutes
            
            per_slide_savings.append({
                'slide': slide['slide'],
                'value': value,
                'unit': 'minutes',
                'raw_text': f"{value_str} {unit}",
                'context': text
            })
        
        # Look for monthly savings
        monthly_matches = re.findall(r'(\d+(?:\.\d+)?)\s*hours?\s*(?:saved\s*)?(?:per\s*consultant\s*monthly)', text, re.IGNORECASE)
        for value_str in monthly_matches:
            monthly_savings.append({
                'slide': slide['slide'],
                'value': float(value_str),
                'unit': 'hours_per_month',
                'raw_text': f"{value_str} hours",
                'context': text
            })
    
    # Check for per-slide conflicts
    if len(set(item['value'] for item in per_slide_savings)) > 1:
        conflicts.append({
            "type": "time_per_slide_conflict",
            "severity": "medium",
            "description": "Found conflicting time savings per slide values",
            "values": per_slide_savings
        })
    
    return conflicts

def detect_sum_breakdown_conflicts(slides):
    """Detect when breakdown components don't sum to the claimed total"""
    conflicts = []
    
    for slide in slides:
        text = slide["text"]
        
        # Look for "X Hours Saved Per Consultant Monthly" pattern
        total_match = re.search(r'(\d+)\s*Hours?\s*Saved\s*Per\s*Consultant\s*Monthly', text, re.IGNORECASE)
        if total_match:
            claimed_total = int(total_match.group(1))
            
            # Look for individual time savings in the same slide
            individual_savings = re.findall(r'(\d+)\s*hours?\s*per\s*consultant\s*monthly', text, re.IGNORECASE)
            if len(individual_savings) > 1:  # More than just the total
                # Remove the total from the list
                individual_values = [int(x) for x in individual_savings if int(x) != claimed_total]
                
                if individual_values:
                    actual_sum = sum(individual_values)
                    if actual_sum != claimed_total:
                        conflicts.append({
                            "type": "sum_breakdown_mismatch",
                            "severity": "high",
                            "slide": slide['slide'],
                            "claimed_total": claimed_total,
                            "breakdown_sum": actual_sum,
                            "breakdown_values": individual_values,
                            "description": f"Claimed total ({claimed_total}) doesn't match sum of breakdown ({actual_sum})"
                        })
    
    return conflicts

def detect_unit_mixing_conflicts(slides):
    """Detect when same metrics are presented in different time units"""
    conflicts = []
    layout_optimization_values = []
    
    for slide in slides:
        text = slide["text"]
        
        # Look for Layout Optimization metrics
        if 'layout optimization' in text.lower():
            # Check for yearly values
            yearly_matches = re.findall(r'(\d+)\s*hours?\s*(?:saved\s*)?(?:per\s*consultant[/\s]*year)', text, re.IGNORECASE)
            for value in yearly_matches:
                layout_optimization_values.append({
                    'slide': slide['slide'],
                    'value': int(value),
                    'unit': 'per_year',
                    'normalized_monthly': int(value) / 12
                })
            
            # Check for monthly values
            monthly_matches = re.findall(r'(\d+)\s*hours?\s*(?:saved\s*)?(?:per\s*consultant\s*monthly)', text, re.IGNORECASE)
            for value in monthly_matches:
                layout_optimization_values.append({
                    'slide': slide['slide'],
                    'value': int(value),
                    'unit': 'per_month',
                    'normalized_monthly': int(value)
                })
    
    # Check if the values are consistent when normalized
    if len(layout_optimization_values) > 1:
        normalized_values = [item['normalized_monthly'] for item in layout_optimization_values]
        if len(set(normalized_values)) > 1:
            conflicts.append({
                "type": "unit_mixing_confusion",
                "severity": "medium",
                "metric": "Layout Optimization",
                "description": "Same metric presented in different time units causing potential confusion",
                "values": layout_optimization_values
            })
    
    return conflicts

# Enhanced numeric conflict detection with context awareness
def detect_contextual_numeric_conflicts(slides):
    """Enhanced version that considers context when detecting numeric conflicts"""
    metric_map = defaultdict(lambda: defaultdict(list))
    conflicts = []

    # Define metric categories to avoid false positives
    metric_categories = {
        'time_savings': ['hours saved', 'minutes saved', 'time saved', 'saves an estimated'],
        'impact': ['$', 'saved in', 'productivity', 'million'],
        'efficiency': ['faster', 'speed', 'efficiency'],
        'comparison': ['vs', 'compared to', 'versus']
    }

    for slide in slides:
        text = slide["text"]
        
        # More sophisticated pattern matching with context
        patterns = [
            r'(?:saves an estimated|delivering|achieving)\s*(\d+(?:\.\d+)?)\s*(hours?|minutes?|mins?)',
            r'(\$\d+(?:\.\d+)?[MK]?)\s*(?:saved|impact)',
            r'(\d+(?:\.\d+)?x)\s*(?:faster|speed)',
            r'(\d+(?:\.\d+)?)\s*(hours?|minutes?)\s*(?:from|saved)',
        ]
        
        for pattern in patterns:
            matches = re.finditer(pattern, text, re.IGNORECASE)
            for match in matches:
                value = match.group(1)
                context = text[max(0, match.start()-50):match.end()+50].strip()
                
                # Categorize the metric based on context
                category = categorize_metric(context)
                if category:
                    metric_key = f"{category}:{value}"
                    metric_map[metric_key][value].append({
                        'slide': slide['slide'],
                        'context': context,
                        'full_match': match.group(0)
                    })

    # Only report conflicts within same categories
    for metric_key, values in metric_map.items():
        if len(values) > 1:
            category = metric_key.split(':')[0]
            conflicts.append({
                "type": "contextual_numeric_conflict",
                "category": category,
                "metric": metric_key,
                "values": dict(values)
            })
    
    return conflicts

def categorize_metric(context):
    """Categorize a metric based on its context"""
    context_lower = context.lower()
    
    if any(word in context_lower for word in ['hours saved', 'minutes saved', 'time saved']):
        return 'time_savings'
    elif any(word in context_lower for word in ['$', 'productivity', 'impact']):
        return 'impact'
    elif any(word in context_lower for word in ['faster', 'speed', 'efficiency']):
        return 'efficiency'
    else:
        return 'other'

# 2. Detecting percentage sum issues (keeping original)
def detect_percent_sum_issues(slides, tolerance=2.0):
    """Checks if percentages in a slide sum to ~100%."""
    issues = []
    for s in slides:
        percents = re.findall(r"[-+]?\d{1,3}(?:\.\d+)?%", s["text"])
        if len(percents) >= 2:
            total = sum(float(p.strip("%")) for p in percents)
            if abs(total - 100.0) > tolerance:
                issues.append({"type": "percent_total_mismatch", "slide": s["slide"], "found_total": total, "details": percents})
    return issues

# Section 4: Enhanced Gemini Analysis

def extract_json_from_text(t):
    """Extract the first JSON array substring from text (crude but practical)."""
    if not isinstance(t, str):
        return None
    start = t.find('[')
    end = t.rfind(']')
    if start != -1 and end != -1 and end > start:
        return t[start:end+1]
    return None

def call_gemini(slide_objs):
    """Enhanced Gemini call with better prompts for specific inconsistency types"""
    slide_texts = []
    for s in slide_objs:
        ppt_excerpt = (s.get("pptx_text") or "")[:1200]
        ocr_excerpt = (s.get("ocr_text") or "")[:1200]
        slide_texts.append(f"SLIDE {s['slide']}:\nPPTX: {ppt_excerpt}\nOCR: {ocr_excerpt}")

    prompt = f"""
    Analyze these presentation slides for factual and logical inconsistencies. Focus on:

    1. MONETARY VALUES: Look for conflicting dollar, or other revenue currency amounts (e.g., "$2M saved" vs "$3M saved")
    2. TIME SAVINGS: Check for inconsistent time values (e.g., "15 mins per slide" vs "20 mins per slide")  
    3. MATHEMATICAL ERRORS: Verify if breakdown components sum to claimed totals
    4. CONTRADICTORY CLAIMS: Find opposing statements about the same topic
    5. TIMELINE CONFLICTS: Identify date or sequence inconsistencies

    Return a JSON array where each element has:
    - issue_type: one of [monetary_conflict, time_savings_conflict, math_error, contradictory_claim, timeline_mismatch, other]
    - severity: [high, medium, low]
    - slides: list of slide numbers involved
    - description: clear explanation of the inconsistency
    - evidence: specific text excerpts that conflict

    If no issues found, return [].

    Slides:
    {chr(10).join(slide_texts)}
    """

    try:
        resp = model.generate_content(prompt)
        text = getattr(resp, "text", None) or str(resp)
    except Exception as e:
        return [{"type": "llm_error", "error": str(e)}]

    json_blob = extract_json_from_text(text)
    if not json_blob:
        try:
            return json.loads(text)
        except Exception:
            return [{"type": "parsing_error", "raw": text}]
    try:
        return json.loads(json_blob)
    except Exception as e:
        return [{"type": "parsing_error", "error": str(e), "raw": json_blob}]

def get_default_paths():
    base_dir = Path(__file__).resolve().parent  # script is in NOOGATASSIGNMENT/
    # the images are expected under NoogatAssignment/NoogatAssignment/
    images_dir = base_dir / "NoogatAssignment" / "NoogatAssignment"
    pptx_files = list(base_dir.rglob("*.pptx"))
    pptx_file = pptx_files[0] if pptx_files else None

    return str(pptx_file) if pptx_file else None, str(images_dir)


def main():
    default_pptx, default_images = get_default_paths()

    print("\n--- ENHANCED NOOGATASSIGNMENT Inconsistency Detector ---")
    print(f"Default PPTX:   {default_pptx if default_pptx else 'Not found'}")
    print(f"Default Images: {default_images if os.path.isdir(default_images) else 'Not found'}")
    print("------------------------------------------------------")

    pptx_path = input(f"Enter PPTX path [{default_pptx}]: ").strip() or default_pptx
    images_path = input(f"Enter images folder [{default_images}]: ").strip() or default_images

    if not pptx_path or not os.path.isfile(pptx_path):
        print(f"❌ PPTX file not found: {pptx_path}")
        return
    if not images_path or not os.path.isdir(images_path):
        print(f"❌ Images folder not found: {images_path}")
        return

    print("\n✅ Paths confirmed:")
    print(f"PPTX:   {pptx_path}")
    print(f"Images: {images_path}")

    # Start processing: extract text from PPTX and images
    print("\n[1/6] Extracting text from PPTX slides...")
    raw_pptx_slides = extract_text_from_pptx(pptx_path)

    print("[2/6] Running OCR on slide images...")
    raw_image_slides = extract_text_from_images(images_path)

    # Normalize both sources
    print("[3/6] Normalizing slides (extracting numbers/dates)...")
    pptx_norm = [normalize_slides(s) for s in tqdm(raw_pptx_slides, desc="Normalizing PPTX")]
    ocr_norm = [normalize_slides(s) for s in tqdm(raw_image_slides, desc="Normalizing OCR")]

    # Merge per-slide (handle differing lengths)
    max_slides = max(len(pptx_norm), len(ocr_norm))
    combined_slides = []
    for i in range(max_slides):
        ppt = pptx_norm[i] if i < len(pptx_norm) else {"slide": i+1, "text": "", "numbers": [], "dates": []}
        ocr = ocr_norm[i] if i < len(ocr_norm) else {"slide": i+1, "text": "", "numbers": [], "dates": []}
        merged_text = (ppt["text"] + " " + ocr["text"]).strip()
        merged_numbers = list(set(ppt.get("numbers", []) + ocr.get("numbers", [])))
        merged_dates = list(set(ppt.get("dates", []) + ocr.get("dates", [])))
        combined_slides.append({
            "slide": i+1,
            "text": merged_text,
            "numbers": merged_numbers,
            "dates": merged_dates,
            "pptx_text": ppt["text"],
            "ocr_text": ocr["text"]
        })

    # Run enhanced rule-based checks
    print("[4/6] Running enhanced rule-based detectors...")
    rule_issues = []
    
    # Original detectors
    rule_issues.extend(detect_percent_sum_issues(combined_slides))
    
    # Enhanced detectors
    rule_issues.extend(detect_impact_value_conflicts(combined_slides))
    rule_issues.extend(detect_time_savings_conflicts(combined_slides))
    rule_issues.extend(detect_sum_breakdown_conflicts(combined_slides))
    rule_issues.extend(detect_unit_mixing_conflicts(combined_slides))
    rule_issues.extend(detect_contextual_numeric_conflicts(combined_slides))

    # LLM-based deep checks (if API key available)
    llm_issues = []
    api_key_present = os.environ.get("my_api_key") is not None
    mode = "deep"

    if mode == "deep":
        if not api_key_present:
            print("⚠️  No Gemini API key found in environment variable 'my_api_key'. Skipping deep LLM checks.")
        else:
            print("[5/6] Running enhanced Gemini deep checks (batched)...")
            batch_size = 8
            for i in range(0, len(combined_slides), batch_size):
                batch = combined_slides[i:i+batch_size]
                res = call_gemini(batch)
                if isinstance(res, list):
                    llm_issues.extend(res)
                else:
                    llm_issues.append({"type":"llm_unexpected", "raw": str(res)})

    # Aggregate and save output with enhanced summary
    print("[6/6] Generating enhanced report...")
    
    # Categorize issues by severity
    high_priority = [issue for issue in rule_issues if issue.get('severity') == 'high']
    medium_priority = [issue for issue in rule_issues if issue.get('severity') == 'medium']
    low_priority = [issue for issue in rule_issues if issue.get('severity') not in ['high', 'medium']]

    out = {
        "summary": {
            "slides_processed": len(combined_slides),
            "total_issues": len(rule_issues) + len(llm_issues),
            "rule_issues": len(rule_issues),
            "llm_issues": len(llm_issues),
            "high_priority_issues": len(high_priority),
            "medium_priority_issues": len(medium_priority),
            "low_priority_issues": len(low_priority)
        },
        "rule_issues": rule_issues,
        "llm_issues": llm_issues,
        "slides": [
            {"slide": s["slide"], "pptx_text": s["pptx_text"], "ocr_text": s["ocr_text"], "numbers": s["numbers"], "dates": s["dates"]}
            for s in combined_slides
        ]
    }
    
    out_path = Path(__file__).resolve().parent / "inconsistencies_enhanced.json"
    with open(out_path, "w", encoding="utf-8") as f:
        json.dump(out, f, indent=2, ensure_ascii=False)

    print("\n=== Enhanced Analysis Complete ===")
    print(f"Slides processed: {out['summary']['slides_processed']}")
    print(f"Total issues found: {out['summary']['total_issues']}")
    print(f"  - High priority: {out['summary']['high_priority_issues']}")
    print(f"  - Medium priority: {out['summary']['medium_priority_issues']}")  
    print(f"  - Low priority: {out['summary']['low_priority_issues']}")
    print(f"Rule-based issues: {out['summary']['rule_issues']}")
    print(f"LLM issues: {out['summary']['llm_issues']}")
    print(f"Saved enhanced report to: {out_path}")

    # Print summary of key findings with enhanced details
    if high_priority:
        print("\n🚨 HIGH PRIORITY ISSUES DETECTED:")
        for issue in high_priority:
            issue_type = issue.get('type', 'Unknown').replace('_', ' ').title()
            description = issue.get('description', 'No description')
            slide_num = issue.get('slide', 'Unknown')
            print(f"  - {issue_type} (Slide {slide_num}): {description}")
    
    if medium_priority:
        print("\n⚠️  MEDIUM PRIORITY ISSUES:")
        for issue in medium_priority:
            issue_type = issue.get('type', 'Unknown').replace('_', ' ').title()
            description = issue.get('description', 'No description')
            slides_involved = issue.get('slides', [issue.get('slide', 'Unknown')])
            if isinstance(slides_involved, list):
                slide_info = f"Slides {', '.join(map(str, slides_involved))}"
            else:
                slide_info = f"Slide {slides_involved}"
            print(f"  - {issue_type} ({slide_info}): {description}")
    
    if llm_issues and any(issue.get('issue_type') != 'llm_error' for issue in llm_issues):
        print("\n🔍 AI-DETECTED ISSUES:")
        for issue in llm_issues:
            if issue.get('issue_type') == 'llm_error':
                continue  # Skip error messages
            issue_type = issue.get('issue_type', 'Unknown').replace('_', ' ').title()
            description = issue.get('description', 'No description')
            slides_involved = issue.get('slides', [])
            if slides_involved:
                slide_info = f"Slides {', '.join(map(str, slides_involved))}"
                print(f"  - {issue_type} ({slide_info}): {description}")
            else:
                print(f"  - {issue_type}: {description}")
    
    if not rule_issues and not llm_issues:
        print("\n✅ NO INCONSISTENCIES DETECTED - All slides appear consistent!")
    
    print(f"\nDetailed report saved to: {out_path}")

if __name__ == "__main__":
    main()